# text_extraction.py
from file_storage import StorageAdapter
from resource_db import ResourceDB
from chunking import chunk_extracted_content
import logging
from typing import List, Dict
from pptx import Presentation
import fitz 
from typing import BinaryIO

logger = logging.getLogger(__name__)

class ResourceProcessor:
    def __init__(self, storage_config: dict, db_path: str = "knowledge_graph/resources.db"):
        self.storage = StorageAdapter(storage_config)
        self.db = ResourceDB(db_path)
    
    def process_new_resources(self):
        pending_resources = self.db.get_pending_resources()
        for resource in pending_resources:
            try:
                self.db.update_processing_status(resource['id'], 'processing')
                
                # Download file from cloud storage
                file_data = self.storage.download_file(resource['file_name'])
                
                # Process based on file type - ADD THIS FUNCTION
                if resource['file_name'].lower().endswith('.pdf'):
                    pages = self.extract_text_from_pdf(file_data)
                elif resource['file_name'].lower().endswith('.pptx'):
                    pages = self.extract_text_from_pptx(file_data)
                else:
                    logger.warning(f"Unsupported file type: {resource['file_name']}")
                    continue
                
                # Prepare for chunking (add source_file reference)
                for page in pages:
                    page['source_file'] = resource['file_name']
                
                # Chunk content
                chunked_content = chunk_extracted_content(
                    {resource['file_name']: pages}, 
                    strategy="semantic"
                )
                chunks = chunked_content.get(resource['file_name'], [])
                
                self.db.save_chunks(resource['id'], chunks)
                self.db.update_processing_status(resource['id'], 'processed')
                logger.info(f"Processed {resource['file_name']} with {len(chunks)} chunks")
                
            except Exception as e:
                logger.error(f"Failed to process {resource['file_name']}: {e}")
                self.db.update_processing_status(resource['id'], 'error')

    def add_new_resource(self, file_name: str, file_type: str, file_data: BinaryIO, metadata: dict = None):
        self.storage.upload_file(file_name, file_data)
        resource_id = self.db.add_resource(file_name, file_type, metadata)
        logger.info(f"Added new resource {file_name} with ID {resource_id}")
        return resource_id

    def extract_text_from_pdf(self, file_data: bytes) -> List[Dict]:
        """Extract text from PDF with metadata"""
        try:
            doc = fitz.open(stream=file_data, filetype="pdf")
            pages = []
            for page_num, page in enumerate(doc, 1):
                text = page.get_text()
                if not text.strip():
                    continue
                    
                # Add metadata
                pages.append({
                    "text": text,
                    "page": page_num,
                    "type": "text",
                    "section": f"Page {page_num}"
                })
            doc.close()
            return pages
        except Exception as e:
            logger.error(f"Error extracting PDF from stream: {e}")
            return []

    # In case got PPTX in the future
    def extract_text_from_pptx(self, file_data: bytes) -> List[Dict]:
        try:
            from io import BytesIO
            prs = Presentation(BytesIO(file_data))
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract slide title
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_num}"
                
                # Extract content with bullet hierarchy
                content = []
                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    if shape == slide.shapes.title:
                        continue
                    
                    # Preserve bullet hierarchy
                    indent_level = 0
                    if shape.text_frame.paragraphs:
                        indent_level = shape.text_frame.paragraphs[0].level
                    content.append({
                        "text": shape.text.strip(),
                        "indent": indent_level,
                        "type": "bullet"
                    })
                
                # Extract speaker notes
                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                
                slides.append({
                    "title": title,
                    "content": content,
                    "notes": notes,
                    "slide": slide_num,
                    "type": "slide",
                    "section": title
                })
            return slides
        except Exception as e:
            logger.error(f"Error extracting PDF from stream: {e}")
            return []

if __name__ == "__main__":
    # Configuration for local development
    config = {
        'storage_type': 'local',
        'base_path': './knowledge_graph/resource'
    }
    
    processor = ResourceProcessor(config)
    
    # Add a new resource (simulating upload)
    with open('new_design.pdf', 'rb') as f:
        processor.add_new_resource('new_design.pdf', 'pdf', f)
    
    # Process all pending resources
    processor.process_new_resources()